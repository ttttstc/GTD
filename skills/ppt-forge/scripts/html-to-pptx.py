#!/usr/bin/env python3
"""
HTML Slide → PPTX 转换脚本

将单个 HTML 文件（包含多个 <section> 标签，每个是一个 slide）
转换为 PPTX 单文件。

依赖：pip install python-pptx beautifulsoup4

用法：
    python html-to-pptx.py input.html output.pptx
    python html-to-pptx.py input.html output.pptx --width 1280 --height 720
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from bs4 import BeautifulSoup
except ImportError:
    print("缺少依赖，请安装：pip install python-pptx beautifulsoup4")
    sys.exit(1)

# 华为品牌色
BRAND_COLORS = {
    'huawei': {
        'primary': RGBColor(0xC7, 0x02, 0x0E),      # 华为红
        'secondary': RGBColor(0xA5, 0x00, 0x00),   # 深红
        'dark': RGBColor(0x1A, 0x1A, 0x1A),        # 黑
        'gray': RGBColor(0xF5, 0xF5, 0xF5),        # 浅灰
        'white': RGBColor(0xFF, 0xFF, 0xFF),
        'text': RGBColor(0x1A, 0x1A, 0x1A),
        'text-secondary': RGBColor(0x66, 0x66, 0x66),
    }
}

# 字号换算：CSS px → PPTX pt
PX_TO_PT = 0.75


def parse_args():
    parser = argparse.ArgumentParser(description='HTML Slide → PPTX')
    parser.add_argument('input', help='输入 HTML 文件路径')
    parser.add_argument('output', help='输出 PPTX 文件路径')
    parser.add_argument('--width', type=int, default=1280, help='幻灯片宽度 (px)')
    parser.add_argument('--height', type=int, default=720, help='幻灯片高度 (px)')
    parser.add_argument('--brand', default='huawei', help='品牌：huawei/apple/ali/generic')
    return parser.parse_args()


def inches(px, total_px):
    """将像素转换为英寸（基于总宽度或高度）"""
    return px / total_px


def add_slide_from_html(prs, section, brand_colors, slide_width, slide_height):
    """从 HTML section 生成一个 PPTX slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 获取 section 内的 HTML 内容
    content = section.decode_contents() if hasattr(section, 'decode_contents') else str(section)
    soup = BeautifulSoup(content, 'html.parser')

    # 解析主要文本元素
    for elem in soup.find_all(['h1', 'h2', 'h3', 'p', 'div', 'span']):
        text = elem.get_text(strip=True)
        if not text:
            continue

        # 估算字号
        font_size_tag = elem.get('style', '')
        px_match = None
        if 'font-size' in font_size_tag:
            import re
            match = re.search(r'font-size:\s*(\d+)px', font_size_tag)
            if match:
                px = int(match.group(1))
                px_match = px

        # 根据标签和样式估算层级
        tag = elem.name.lower()
        if tag == 'h1':
            font_size = int(px_match) if px_match else 26
            bold = True
        elif tag == 'h2':
            font_size = int(px_match) if px_match else 14
            bold = True
        elif tag == 'h3':
            font_size = int(px_match) if px_match else 12
            bold = True
        else:
            font_size = int(px_match) if px_match else 10
            bold = 'font-weight' in font_size_tag and '700' in font_size_tag

        # 获取颜色
        color = brand_colors['text']
        if 'color' in font_size_tag:
            # 简化处理，实际应解析 hex 颜色
            pass

        # 获取位置（简化版：按文档顺序堆叠）
        # 实际应解析 CSS 布局
        left = Inches(0.5)
        top = Inches(sum([s.height for s in prs.slides[:-1]]) / 914400 or 0.5)
        width = Inches(slide_width / 96 * 0.9)
        height = Inches(font_size * PX_TO_PT / 72)

        # 添加文本框
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text[:200]  # 截断超长文本
        p.font.size = Pt(font_size * PX_TO_PT)
        p.font.bold = bold
        p.font.color.rgb = color

    return slide


def main():
    args = parse_args()

    # 读取 HTML
    with open(args.input, 'r', encoding='utf-8') as f:
        html_content = f.read()

    # 解析 slides
    soup = BeautifulSoup(html_content, 'html.parser')
    sections = soup.find_all('section')

    if not sections:
        print("错误：未找到 <section> 标签，请确保 HTML 包含多个 <section> 元素")
        sys.exit(1)

    # 创建 PPT
    prs = Presentation()
    prs.slide_width = Inches(args.width / 96)
    prs.slide_height = Inches(args.height / 96)

    brand = args.brand.lower()
    colors = BRAND_COLORS.get(brand, BRAND_COLORS['huawei'])

    # 转换每个 section
    for i, section in enumerate(sections):
        print(f"转换 slide {i + 1}/{len(sections)}...")
        add_slide_from_html(prs, section, colors, args.width, args.height)

    # 保存
    output_path = Path(args.output)
    prs.save(str(output_path))
    print(f"完成：{output_path} ({len(sections)} slides)")


if __name__ == '__main__':
    main()
