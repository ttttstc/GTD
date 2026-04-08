#!/usr/bin/env python3
"""
ZaraPPT PPTX 提取脚本

从 PPTX 文件中提取文本内容并转换为 Markdown 格式。

依赖：pip install python-pptx

用法：
    python zarappt-extract-pptx.py input.pptx
    python zarappt-extract-pptx.py input.pptx -o output.md
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    print("缺少依赖，请安装：pip install python-pptx")
    sys.exit(1)


def parse_args():
    parser = argparse.ArgumentParser(description='PPTX 内容提取')
    parser.add_argument('input', help='输入 PPTX 文件路径')
    parser.add_argument('-o', '--output', help='输出 Markdown 文件路径')
    return parser.parse_args()


def extract_pptx(pptx_path):
    """从 PPTX 提取文本内容"""
    prs = Presentation(pptx_path)

    slides_content = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        slide_title = None

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    # 尝试识别标题（第一个大字号文本）
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if paragraph.runs:
                                first_run = paragraph.runs[0]
                                if first_run.font.size and first_run.font.size > Pt(24):
                                    slide_title = text
                                    break

                    if not slide_title or text != slide_title:
                        slide_text.append(text)

        slide_content = {
            'number': i,
            'title': slide_title or f'Slide {i}',
            'content': '\n'.join(slide_text) if slide_text else '(无文本内容)'
        }
        slides_content.append(slide_content)

    return slides_content


def format_markdown(slides_content):
    """格式化为 Markdown"""
    md_lines = ['# PPT 内容提取\n']

    for slide in slides_content:
        md_lines.append(f'## Slide {slide["number"]}: {slide["title"]}\n')
        md_lines.append(slide['content'])
        md_lines.append('\n---\n')

    return '\n'.join(md_lines)


def main():
    args = parse_args()

    # 检查输入文件
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"错误：文件不存在 {args.input}")
        sys.exit(1)

    # 提取内容
    print(f"正在提取: {args.input}")
    slides_content = extract_pptx(input_path)

    # 格式化
    markdown = format_markdown(slides_content)

    # 输出
    if args.output:
        output_path = Path(args.output)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(markdown)
        print(f"已保存到: {args.output}")
    else:
        print(markdown)

    print(f"\n共提取 {len(slides_content)} 张幻灯片")


if __name__ == '__main__':
    main()
